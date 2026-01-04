#!/usr/bin/env python3

from __future__ import annotations

import json
import os
import re
import sys
import time
import hashlib
import multiprocessing as mp
from datetime import datetime
from pathlib import Path
from queue import Queue, Empty
from threading import Event, Lock, Thread
from typing import Dict, List, Set, Tuple
from urllib.parse import parse_qsl, urlencode, urljoin, urlparse, urldefrag, urlunparse

import requests
from bs4 import BeautifulSoup

PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from app.config import Config

RAW_DIR = PROJECT_ROOT / "app" / "kb" / "raw"
INDEX_DIR = PROJECT_ROOT / "index"
INDEX_FILE = INDEX_DIR / "index.json"


def log(level: str, message: str) -> None:
    ts = datetime.utcnow().isoformat(timespec="seconds")
    print(f"{ts} [{level}] {message}", flush=True)


def _env(name: str, default: str) -> str:
    return os.getenv(name, default).strip()


REQUEST_TIMEOUT = int(_env("REQUEST_TIMEOUT", "20") or "20")
CRAWL_DELAY_SECONDS = float(_env("CRAWL_DELAY_SECONDS", "0.25") or "0.25")
CRAWL_WORKERS = int(_env("CRAWL_WORKERS", "12") or "12")
MAX_PAGES_TOTAL = int(_env("MAX_PAGES_TOTAL", "20000") or "20000")
MAX_PAGES_PER_DOMAIN = int(_env("MAX_PAGES_PER_DOMAIN", "2000") or "2000")
PERSIST_EVERY = int(_env("PERSIST_EVERY", "200") or "200")
USER_AGENT = _env("USER_AGENT", "CyberTrainingMegaIndexer/3.0")

CHUNK_WORDS = int(_env("CHUNK_WORDS", "220") or "220")
CHUNK_OVERLAP = int(_env("CHUNK_OVERLAP", "40") or "40")

SKIP_PDF = _env("SKIP_PDF", "0") == "1"
PDF_MAX_PAGES = int(_env("PDF_MAX_PAGES", "60") or "60")
PDF_TIMEOUT_SECONDS = int(_env("PDF_TIMEOUT_SECONDS", "12") or "12")

SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".csv", ".xlsx"}

RELEVANCE_KEYWORDS = {
    "vulnerability",
    "vulnerabilities",
    "exploit",
    "malware",
    "phishing",
    "ransomware",
    "cve",
    "nist",
    "iso 27001",
    "owasp",
    "incident response",
    "incident",
    "threat",
    "attack",
    "mitigation",
    "encryption",
    "authentication",
    "authorization",
    "risk",
    "compliance",
    "soc",
    "siem",
    "zero trust",
    "firewall",
    "intrusion",
    "patching",
}
KEYWORD_THRESHOLD = int(_env("KEYWORD_THRESHOLD", "2") or "2")
MIN_WORDS_FOR_LOW_SIGNAL = int(_env("MIN_WORDS_FOR_LOW_SIGNAL", "80") or "80")

PATH_ALLOW_PATTERNS = [
    "/guidance",
    "/publications",
    "/alerts",
    "/advisories",
    "/standards",
    "/resources",
    "/learn",
]
PATH_DENY_PATTERNS = [
    "/about",
    "/careers",
    "/jobs",
    "/events",
    "/contact",
    "/press",
    "/newsroom",
    "/privacy",
]

try:
    from PyPDF2 import PdfReader  # type: ignore
except Exception:
    PdfReader = None

try:
    from docx import Document  # type: ignore
except Exception:
    Document = None

try:
    from pptx import Presentation  # type: ignore
except Exception:
    Presentation = None

try:
    import openpyxl  # type: ignore
except Exception:
    openpyxl = None


def strip_www(host: str) -> str:
    host = (host or "").lower().strip()
    return host[4:] if host.startswith("www.") else host


def canonicalize_url(url: str) -> str:
    url, _ = urldefrag(url)
    parsed = urlparse(url)
    scheme = parsed.scheme or "https"
    host = strip_www(parsed.netloc)
    query = [
        (k, v)
        for k, v in parse_qsl(parsed.query, keep_blank_values=True)
        if not (k.lower().startswith("utm_") or k.lower() in {"gclid", "fbclid"})
    ]
    path = parsed.path or "/"
    if path != "/" and path.endswith("/"):
        path = path[:-1]
    return urlunparse((scheme, host, path, "", urlencode(query), ""))


def gather_allowed_domains() -> Tuple[List[str], List[str]]:
    raw_domains = Config.allowed_domains or []
    normalized_domains: List[str] = []
    roots: List[str] = []
    seen_hosts: Set[str] = set()
    for entry in raw_domains:
        entry = (entry or "").strip()
        if not entry:
            continue
        if entry.startswith("http://") or entry.startswith("https://"):
            url = entry
        else:
            url = f"https://{entry}"
        canon = canonicalize_url(url)
        host = strip_www(urlparse(canon).netloc)
        if not host or host in seen_hosts:
            continue
        seen_hosts.add(host)
        normalized_domains.append(host)
        roots.append(f"https://{host}")
    return normalized_domains, roots


ALLOWED_DOMAINS, SEED_URLS = gather_allowed_domains()
log("INFO", f"Allowed domains (from Config.allowed_domains): {ALLOWED_DOMAINS}")
log("INFO", f"Crawl roots derived from allowed domains: {SEED_URLS}")


def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def sha256_hex(text: str) -> str:
    return hashlib.sha256((text or "").encode("utf-8", errors="ignore")).hexdigest()


def tokenize(text: str) -> Set[str]:
    return set(re.findall(r"[a-z0-9]+", (text or "").lower()))


def chunk_words(text: str, size: int, overlap: int) -> List[Tuple[str, str]]:
    words = (text or "").split()
    if not words:
        return []
    chunks: List[Tuple[str, str]] = []
    start = 0
    while start < len(words):
        end = min(len(words), start + size)
        chunk = " ".join(words[start:end]).strip()
        if chunk:
            chunks.append((f"words={start}-{end}", chunk))
        if end == len(words):
            break
        start = max(end - overlap, start + 1)
    return chunks


def relevance_score(text: str) -> Tuple[int, int, bool]:
    lower = text.lower()
    keyword_hits = sum(1 for kw in RELEVANCE_KEYWORDS if kw in lower)
    word_count = len(text.split())
    if keyword_hits >= KEYWORD_THRESHOLD:
        return keyword_hits, word_count, True
    if keyword_hits > 0 and word_count >= MIN_WORDS_FOR_LOW_SIGNAL:
        return keyword_hits, word_count, True
    return keyword_hits, word_count, False


def path_allowed(url: str, parent_relevant: bool) -> bool:
    parsed = urlparse(url)
    path = (parsed.path or "/").lower()
    for pattern in PATH_DENY_PATTERNS:
        if pattern and pattern in path:
            if parent_relevant:
                return True
            log("INFO", f"skipped_path_filter (deny): {url}")
            return False
    if PATH_ALLOW_PATTERNS and not parent_relevant:
        if not any(pattern in path for pattern in PATH_ALLOW_PATTERNS):
            if path.count("/") > 4:
                log("INFO", f"skipped_path_filter (allow preference): {url}")
                return False
    return True


def persist_index(documents: List[Dict]) -> None:
    payload = {"trusted_domains": ALLOWED_DOMAINS, "documents": documents}
    INDEX_DIR.mkdir(parents=True, exist_ok=True)
    tmp_path = INDEX_FILE.with_suffix(".tmp")
    with open(tmp_path, "w", encoding="utf-8") as handle:
        json.dump(payload, handle, indent=2, ensure_ascii=False)
    tmp_path.replace(INDEX_FILE)


def maybe_persist(documents: List[Dict]) -> None:
    if PERSIST_EVERY <= 0:
        return
    if len(documents) % PERSIST_EVERY == 0:
        persist_index(documents)


def is_allowed_url(url: str) -> bool:
    try:
        host = strip_www(urlparse(url).netloc)
        if not host:
            return False
        return host in ALLOWED_DOMAINS or any(
            host.endswith(f".{domain}") for domain in ALLOWED_DOMAINS
        )
    except Exception:
        return False


def _pdf_worker(path: str, max_pages: int, q: mp.Queue) -> None:
    try:
        if PdfReader is None:
            q.put([])
            return
        reader = PdfReader(path)
        out: List[Tuple[str, str]] = []
        total = len(reader.pages)
        limit = min(total, max_pages) if max_pages > 0 else total
        for idx in range(limit):
            try:
                text = clean_text(reader.pages[idx].extract_text() or "")
                if text:
                    out.append((f"page={idx+1}", text))
            except Exception:
                continue
        q.put(out)
    except Exception:
        q.put([])


def extract_pdf_safe(path: Path, max_pages: int, timeout: int) -> List[Tuple[str, str]]:
    if PdfReader is None:
        log("WARNING", f"Skipping PDF (PyPDF2 missing): {path.name}")
        return []
    queue: mp.Queue = mp.Queue()
    proc = mp.Process(
        target=_pdf_worker, args=(str(path), max_pages, queue), daemon=True
    )
    proc.start()
    proc.join(timeout)
    if proc.is_alive():
        proc.terminate()
        proc.join(2)
        log("WARNING", f"PDF timeout, skipped: {path.name}")
        return []
    try:
        return queue.get_nowait()
    except Exception:
        return []


def extract_file_segments(path: Path) -> List[Tuple[str, str]]:
    ext = path.suffix.lower()
    if ext in {".txt", ".md", ".csv"}:
        text = clean_text(path.read_text(encoding="utf-8", errors="ignore"))
        return [("file", text)] if text else []
    if ext == ".pdf":
        if SKIP_PDF:
            log("INFO", f"Skipping PDF by configuration: {path.name}")
            return []
        log(
            "INFO",
            f"Parsing PDF: {path.name} (max_pages={PDF_MAX_PAGES}, timeout={PDF_TIMEOUT_SECONDS}s)",
        )
        return extract_pdf_safe(path, PDF_MAX_PAGES, PDF_TIMEOUT_SECONDS)
    if ext == ".docx":
        if not Document:
            log("WARNING", f"python-docx missing; skipping {path.name}")
            return []
        doc = Document(str(path))
        text = clean_text(
            " ".join(p.text for p in doc.paragraphs if p.text and p.text.strip())
        )
        return [("file", text)] if text else []
    if ext == ".pptx":
        if not Presentation:
            log("WARNING", f"python-pptx missing; skipping {path.name}")
            return []
        slides: List[Tuple[str, str]] = []
        prs = Presentation(str(path))
        for idx, slide in enumerate(prs.slides, start=1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text and shape.text.strip()
            ]
            text = clean_text(" ".join(texts))
            if text:
                slides.append((f"slide={idx}", text))
        return slides
    if ext == ".xlsx":
        if not openpyxl:
            log("WARNING", f"openpyxl missing; skipping {path.name}")
            return []
        workbook = openpyxl.load_workbook(path, data_only=True)
        segments = []
        for sheet in workbook.worksheets:
            rows = [
                " ".join(str(cell) for cell in row if cell is not None).strip()
                for row in sheet.iter_rows(values_only=True)
            ]
            text = clean_text(" ".join(r for r in rows if r))
            if text:
                segments.append((f"sheet={sheet.title}", text))
        return segments
    return []


def crawl_websites(max_pages: int) -> List[Dict]:
    if not SEED_URLS:
        log("WARNING", "No crawl roots available; add domains to ALLOWED_DOMAINS.")
        return []
    task_queue: Queue[str] = Queue()
    visited: Set[str] = set()
    scheduled: Set[str] = set()
    pages: List[Dict] = []
    domain_counts: Dict[str, int] = {}
    state_lock = Lock()
    stop_event = Event()

    for seed in SEED_URLS:
        canonical = canonicalize_url(seed)
        host = strip_www(urlparse(canonical).netloc)
        if not host:
            log("WARNING", f"Invalid seed: {seed}")
            continue
        with state_lock:
            if canonical in scheduled:
                continue
            scheduled.add(canonical)
        task_queue.put(canonical)
        log("INFO", f"Seed queued: {canonical}")

    log(
        "INFO",
        f"Beginning crawl with {len(scheduled)} roots and {CRAWL_WORKERS} workers",
    )

    def worker(worker_id: int) -> None:
        session = requests.Session()
        session.headers.update({"User-Agent": USER_AGENT})
        while True:
            if stop_event.is_set() and task_queue.empty():
                return
            try:
                current = task_queue.get(timeout=1)
            except Empty:
                if stop_event.is_set():
                    return
                continue
            with state_lock:
                if current in visited:
                    task_queue.task_done()
                    continue
                visited.add(current)
            log(
                "INFO",
                f"[worker={worker_id}] Visiting {current} "
                f"(visited={len(visited)} indexed={len(pages)} queue={task_queue.qsize()})",
            )
            host = strip_www(urlparse(current).netloc)
            try:
                resp = session.get(current, timeout=REQUEST_TIMEOUT)
                resp.raise_for_status()
            except Exception as exc:
                log(
                    "WARNING", f"[worker={worker_id}] Request failed: {current} ({exc})"
                )
                task_queue.task_done()
                continue
            soup = BeautifulSoup(resp.text, "html.parser")
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()
            container = soup.find("main") or soup.find("article") or soup.body
            if not container:
                task_queue.task_done()
                continue
            text = clean_text(container.get_text(" "))
            page_added = False
            page_relevant = False
            keyword_hits = 0
            word_count = 0
            if text:
                keyword_hits, word_count, page_relevant = relevance_score(text)
                if page_relevant:
                    with state_lock:
                        domain_counts.setdefault(host, 0)
                        if domain_counts[host] >= MAX_PAGES_PER_DOMAIN:
                            log(
                                "INFO",
                                f"[worker={worker_id}] Domain limit reached for {host}",
                            )
                        elif len(pages) < max_pages:
                            pages.append(
                                {
                                    "id": f"web_{len(pages)}",
                                    "title": (
                                        clean_text(soup.title.string)
                                        if soup.title and soup.title.string
                                        else current
                                    ),
                                    "url": current,
                                    "text": text,
                                }
                            )
                            domain_counts[host] += 1
                            page_added = True
                            log(
                                "INFO",
                                f"[worker={worker_id}] indexed_relevant "
                                f"(keywords={keyword_hits} words={word_count}): {current}",
                            )
                        else:
                            stop_event.set()
                else:
                    log(
                        "INFO",
                        f"[worker={worker_id}] skipped_low_relevance "
                        f"(keywords={keyword_hits} words={word_count}): {current}",
                    )
            discovered = 0
            enqueued = 0
            if not stop_event.is_set():
                for anchor in soup.find_all("a", href=True):
                    href = (anchor.get("href") or "").strip()
                    if not href or href.startswith(("mailto:", "javascript:", "tel:")):
                        continue
                    nxt = canonicalize_url(urljoin(current, href))
                    if not nxt or not is_allowed_url(nxt):
                        continue
                    discovered += 1
                    next_host = strip_www(urlparse(nxt).netloc)
                    if not path_allowed(nxt, page_relevant):
                        continue
                    should_enqueue = False
                    with state_lock:
                        domain_counts.setdefault(next_host, 0)
                        if (
                            domain_counts[next_host] < MAX_PAGES_PER_DOMAIN
                            and nxt not in scheduled
                        ):
                            scheduled.add(nxt)
                            should_enqueue = True
                    if should_enqueue:
                        task_queue.put(nxt)
                        enqueued += 1
            log(
                "INFO",
                f"[worker={worker_id}] discovered={discovered} enqueued={enqueued} queue={task_queue.qsize()}",
            )
            task_queue.task_done()
            if len(pages) >= max_pages:
                stop_event.set()
            if CRAWL_DELAY_SECONDS > 0:
                time.sleep(CRAWL_DELAY_SECONDS)

    threads: List[Thread] = []
    for idx in range(max(1, CRAWL_WORKERS)):
        thread = Thread(target=worker, args=(idx + 1,), daemon=True)
        thread.start()
        threads.append(thread)

    task_queue.join()
    stop_event.set()
    for thread in threads:
        thread.join(timeout=2)

    log("INFO", f"Crawl finished. Visited={len(visited)} Indexed={len(pages)}")
    return pages


def build_index() -> Dict:
    documents: List[Dict] = []
    pages = crawl_websites(MAX_PAGES_TOTAL)
    for page in pages:
        for span, chunk in chunk_words(page["text"], CHUNK_WORDS, CHUNK_OVERLAP):
            documents.append(
                {
                    "id": f"{page['id']}_{span}",
                    "title": page["title"],
                    "content": chunk,
                    "bullets": [],
                    "source": page["url"],
                    "citations": [{"type": "url", "url": page["url"], "locator": span}],
                    "keywords": sorted(tokenize(chunk)),
                    "trusted": True,
                    "meta": {"content_hash": sha256_hex(chunk), "origin": "web"},
                }
            )
            maybe_persist(documents)
    if RAW_DIR.exists():
        files = [
            path
            for path in RAW_DIR.rglob("*")
            if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
        ]
        log("INFO", f"Processing {len(files)} local files from {RAW_DIR}")
        for idx, path in enumerate(files, start=1):
            log(
                "INFO",
                f"[{idx}/{len(files)}] Local file: {path.relative_to(PROJECT_ROOT)}",
            )
            segments = extract_file_segments(path)
            if not segments:
                continue
            rel_source = path.relative_to(PROJECT_ROOT).as_posix()
            for locator, seg_text in segments:
                for span, chunk in chunk_words(seg_text, CHUNK_WORDS, CHUNK_OVERLAP):
                    documents.append(
                        {
                            "id": f"file_{path.stem}_{span}",
                            "title": path.stem,
                            "content": chunk,
                            "bullets": [],
                            "source": rel_source,
                            "citations": [
                                {"type": "file", "path": rel_source, "locator": locator}
                            ],
                            "keywords": sorted(tokenize(chunk)),
                            "trusted": True,
                            "meta": {
                                "content_hash": sha256_hex(chunk),
                                "origin": "local",
                            },
                        }
                    )
                    maybe_persist(documents)
    else:
        log("WARNING", f"Local KB directory not found: {RAW_DIR}")
    persist_index(documents)
    return {"trusted_domains": ALLOWED_DOMAINS, "documents": documents}


def main() -> None:
    log("INFO", "Mega indexer started")
    if not ALLOWED_DOMAINS:
        log(
            "WARNING",
            "Config.allowed_domains is empty. Populate ALLOWED_DOMAINS to crawl.",
        )
    INDEX_DIR.mkdir(parents=True, exist_ok=True)
    index = build_index()
    log("INFO", f"Index built successfully. Documents: {len(index['documents'])}")
    log("INFO", f"Output: {INDEX_FILE}")


if __name__ == "__main__":
    mp.freeze_support()
    main()
