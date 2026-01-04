#!/usr/bin/env python3
"""
Builds an index from heterogeneous knowledge base files.

Reads every supported document under kb/raw and emits JSON in index/index.json:
{
    "trusted_domains": [...],
    "documents": [
        {
            "id": "...",
            "title": "...",
            "content": "...",
            "bullets": [],
            "source": "...",
            "keywords": [...],
            "trusted": true
        }
    ]
}
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Iterable, List

from config import Config

try:
    from PyPDF2 import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None

try:
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

try:
    import openpyxl
except ImportError:  # pragma: no cover
    openpyxl = None


RAW_DIR = Path("app/kb/raw")
INDEX_PATH = Path("index/index.json")
SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".pptx",
    ".csv",
    ".xlsx",
}


def clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def tokenize(text: str) -> List[str]:
    return re.findall(r"[a-z0-9]+", text.lower())


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 120) -> Iterable[str]:
    if not text:
        return []
    tokens = text.split()
    if not tokens:
        return []
    start = 0
    chunks = []
    while start < len(tokens):
        end = min(len(tokens), start + chunk_size)
        chunk = " ".join(tokens[start:end])
        chunks.append(chunk)
        if end == len(tokens):
            break
        start = max(end - overlap, start + 1)
    return chunks


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()

    if suffix in {".txt", ".md"}:
        return clean_text(path.read_text(encoding="utf-8", errors="ignore"))

    if suffix == ".csv":
        return clean_text(path.read_text(encoding="utf-8", errors="ignore"))

    if suffix == ".pdf":
        if not PdfReader:
            raise RuntimeError("PyPDF2 is required for PDF support")
        reader = PdfReader(str(path))
        pages = [page.extract_text() or "" for page in reader.pages]
        return clean_text(" ".join(pages))

    if suffix == ".docx":
        if not Document:
            raise RuntimeError("python-docx is required for DOCX support")
        doc = Document(str(path))
        paragraphs = [para.text for para in doc.paragraphs]
        return clean_text(" ".join(paragraphs))

    if suffix == ".pptx":
        if not Presentation:
            raise RuntimeError("python-pptx is required for PPTX support")
        prs = Presentation(str(path))
        slides = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides.append(shape.text)
        return clean_text(" ".join(slides))

    if suffix == ".xlsx":
        if not openpyxl:
            raise RuntimeError("openpyxl is required for XLSX support")
        book = openpyxl.load_workbook(path, data_only=True)
        segments = []
        for sheet in book.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text:
                    segments.append(row_text)
        return clean_text(" ".join(segments))

    raise ValueError(f"Unsupported file type: {path}")


def build_entries() -> List[dict]:
    documents: List[dict] = []
    if not RAW_DIR.exists():
        return documents

    for path in RAW_DIR.rglob("*"):
        if not path.is_file():
            continue
        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue

        text = extract_text(path)
        if not text:
            continue

        chunks = chunk_text(text)
        rel_source = str(path.relative_to(RAW_DIR))
        base_id = re.sub(r"[^a-z0-9]+", "_", path.stem.lower())
        for idx, chunk in enumerate(chunks, start=1):
            doc_id = f"{base_id}_{idx}"
            documents.append(
                {
                    "id": doc_id,
                    "title": f"{path.stem} (section {idx})",
                    "content": chunk,
                    "bullets": [],
                    "source": f"kb/raw/{rel_source}",
                    "keywords": list(set(tokenize(chunk))),
                    "trusted": True,
                }
            )
    return documents


def main():
    parser = argparse.ArgumentParser(description="Build KB index from local files.")
    parser.add_argument(
        "--output",
        type=Path,
        default=INDEX_PATH,
        help="Output path for index JSON (default: index/index.json)",
    )
    args = parser.parse_args()

    documents = build_entries()
    args.output.parent.mkdir(parents=True, exist_ok=True)
    payload = {
        "trusted_domains": Config.allowed_domains,
        "documents": documents,
    }
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(payload, f, indent=2, ensure_ascii=False)
    print(f"Indexed {len(documents)} document chunks into {args.output}")


if __name__ == "__main__":
    main()
